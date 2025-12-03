"""
Document Processor Module
Handles PDF/PowerPoint/Word/Video processing and Gemini-powered Q&A
"""

import os
import PyPDF2
from pptx import Presentation
from docx import Document
import google.generativeai as genai
from moviepy.editor import VideoFileClip
from openai import OpenAI
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Initialize OpenAI client for Whisper
openai_client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

def extract_text_from_pdf(filepath):
    """Extract text from a PDF file"""
    try:
        text = ""
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                text += f"\n--- Page {page_num + 1} ---\n"
                text += page_text
        
        return text, num_pages
    except Exception as e:
        return f"Error extracting PDF: {str(e)}", 0

def extract_text_from_pptx(filepath):
    """Extract text from a PowerPoint file"""
    try:
        text = ""
        prs = Presentation(filepath)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            text += f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text, len(prs.slides)
    except Exception as e:
        return f"Error extracting PowerPoint: {str(e)}", 0

def extract_text_from_docx(filepath):
    """Extract text from a Word document"""
    try:
        doc = Document(filepath)
        text = ""
        
        for para_num, para in enumerate(doc.paragraphs, start=1):
            if para.text.strip():
                text += para.text + "\n"
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text += " | ".join(row_text) + "\n"
        
        # Count paragraphs as "pages" approximation
        num_pages = max(1, len(doc.paragraphs) // 20)  # Rough estimate
        
        return text, num_pages
    except Exception as e:
        return f"Error extracting Word document: {str(e)}", 0

def extract_text_from_video(filepath):
    """
    Extract audio from video and transcribe using Gemini Audio API
    """
    try:
        # Load video
        video = VideoFileClip(filepath)
        duration_minutes = video.duration / 60
        
        # Check if video has audio
        if video.audio is None:
            video.close()
            return f"Video has no audio track. Cannot transcribe.", int(duration_minutes)
        
        # Save audio temporarily
        audio_path = filepath.rsplit('.', 1)[0] + '_audio.mp3'
        
        try:
            video.audio.write_audiofile(
                audio_path, 
                verbose=False, 
                logger=None,
                codec='mp3'
            )
        except Exception as e:
            video.close()
            return f"Error extracting audio: {str(e)}", int(duration_minutes)
        
        video.close()
        
        # Upload audio to Gemini for transcription
        try:
            audio_file = genai.upload_file(path=audio_path)
            
            # Use Gemini to transcribe
            model = genai.GenerativeModel('gemini-1.5-flash')
            response = model.generate_content([
                "Transcribe this video audio exactly as spoken. Provide a clear, accurate transcription of all speech.",
                audio_file
            ])
            
            transcript = response.text
            
            # Clean up
            genai.delete_file(audio_file.name)
            
        except Exception as e:
            # Clean up audio file even if transcription fails
            if os.path.exists(audio_path):
                os.remove(audio_path)
            return f"Error transcribing audio with Gemini: {str(e)}", int(duration_minutes)
        
        # Clean up temporary audio file
        if os.path.exists(audio_path):
            os.remove(audio_path)
        
        # Format transcript
        text = f"""
--- VIDEO TRANSCRIPT ---
Duration: {duration_minutes:.1f} minutes

{transcript}

--- END TRANSCRIPT ---
"""
        
        return text, int(duration_minutes)
        
    except Exception as e:
        return f"Error processing video: {str(e)}", 0

def process_document(filepath):
    """
    Process an uploaded document (PDF, PowerPoint, Word, or Video)
    Extracts text content from documents or transcribes audio from videos
    """
    filename = os.path.basename(filepath)
    file_ext = os.path.splitext(filename)[1].lower()
    
    # Extract text based on file type
    if file_ext == '.pdf':
        content, num_pages = extract_text_from_pdf(filepath)
        file_type = "pdf"
        unit = "pages"
    elif file_ext in ['.pptx', '.ppt']:
        content, num_pages = extract_text_from_pptx(filepath)
        file_type = "pptx"
        unit = "slides"
    elif file_ext in ['.docx', '.doc']:
        content, num_pages = extract_text_from_docx(filepath)
        file_type = "docx"
        unit = "pages"
    elif file_ext in ['.mp4', '.mov', '.avi', '.mkv', '.webm']:
        content, num_pages = extract_text_from_video(filepath)
        file_type = "video"
        unit = "minutes"
    else:
        return {
            "filename": filename,
            "content": "Unsupported file type",
            "metadata": {"error": f"File type {file_ext} not supported. Use PDF, DOCX, PPTX, or video files."}
        }
    
    # Create document info
    doc_info = {
        "filename": filename,
        "content": content,
        "metadata": {
            "pages": num_pages,
            "file_type": file_type,
            "char_count": len(content),
            "unit": unit
        }
    }
    
    return doc_info

def answer_question_with_gemini(question, documents):
    """
    Answer a question using Google Gemini based on uploaded documents
    """
    
    if not documents:
        return "No documents uploaded yet. Please upload some training materials first!"
    
    try:
        # Combine all document content
        combined_content = ""
        for doc in documents:
            combined_content += f"\n\n=== Content from {doc['filename']} ===\n"
            combined_content += doc['content']
        
        # Create the prompt for Gemini
        prompt = f"""You are a helpful training assistant. Answer the following question based on the provided training materials.

TRAINING MATERIALS:
{combined_content[:15000]}  

QUESTION: {question}

Please provide a clear, detailed answer based on the materials above. If the information isn't in the materials, say so. Include relevant quotes or references to specific sections when appropriate.

ANSWER:"""
        
        # Call Gemini API (using latest model)
        model = genai.GenerativeModel('gemini-2.5-flash')
        response = model.generate_content(prompt)
        
        # Format the response
        answer = response.text
        
        # Add source information
        sources = [doc['filename'] for doc in documents]
        formatted_response = f"""**Answer:**

{answer}

---

**📚 Sources:** {', '.join(sources)}

*Powered by Google Gemini*
"""
        
        return formatted_response
        
    except Exception as e:
        return f"Error generating answer: {str(e)}\n\nPlease check your Gemini API key in the .env file."
